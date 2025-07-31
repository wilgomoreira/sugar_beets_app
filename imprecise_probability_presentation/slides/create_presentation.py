from pptx import Presentation

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Imprecise Probability"
slide.placeholders[1].text = "Overview of the Imprecise Probability\nAuthor: [Wilgo Nunes]\nDate: [06/2025]"

# Slide 2: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Introduction"
slide.placeholders[1].text = (
    "Brief introduction to the concept of imprecise probabilities.\n"
    "Importance in machine learning and uncertainty quantification."
)

# Slide 3: Loading Data
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Loading Data"
slide.placeholders[1].text = (
    "Description of the data loading process.\n"
    "Mention of calibration and test logits/likelihoods.\n"
)

# Slide 4: Obtaining Lower Bound
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Obtaining Lower Bound"
slide.placeholders[1].text = (
    "Lower bound incorporates entropy to penalize uncertain predictions.\n"
    "Formula:\n"
    "lower_bound = \\max\\left(0, \\min\\left(1, \\text{probs} - \\alpha \\times \\text{entropy}\\right)\\right)\n"
    "Alpha controls the strength of the penalty."
)

# Slide 5: Visualizing Likelihood Distributions
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Visualizing Likelihood Distributions"
slide.placeholders[1].text = (
    "KDE plots show the distribution of lower bounds for each class.\n"
    "Helps understand class-wise confidence.\n"
)

# Slide 6: Visualizing Likelihood Distributions
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Visualizing Likelihood Distributions"
slide.placeholders[1].text = (
    "KDE plots show the distribution of lower bounds for each class.\n"
    "Helps understand class-wise confidence.\n"
)

# Slide 7: Metrics for Evaluation
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Metrics for Evaluation"
slide.placeholders[1].text = (
    "Average Precision (AP) & Mean AP (mAP): Measures segmentation quality\n"
    "Expected Calibration Error (ECE): Measures how well probabilities reflect true likelihood.\n"
    "Other metrics: Brier Score, NLL, Entropy (can be added)\n"
)

# Slide 8: Threshold Selection via Quartiles
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Threshold Selection via Quartiles"
slide.placeholders[1].text = (
    "Thresholds for each class are set using quantiles of lower bound distributions.\n"
    "Optimized using cross-validation and parallel grid search.\n"
)

# Slide 9: Superpixel-Based Uncertainty
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Superpixel-Based Uncertainty"
slide.placeholders[1].text = (
    "Superpixels group pixels for region-based uncertainty.\n"
    "Uncertain superpixels are identified and handled differently..\n"
)

# Slide 10: Smooth Interpolation
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Smooth Interpolation"
slide.placeholders[1].text = (
    "Blends confident and baseline probabilities using a parameter alpha.\n"
    "Allows for a soft transition between confident and uncertain predictions.\n"
)

# Slide 11: Results & Metrics
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Results & Metrics"
slide.placeholders[1].text = (
    "Baseline vs. Imprecise Probability metrics (mAP, ECE).\n"
    "Visual comparison of predictions and uncertainty.\n"
)

# Slide 12: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
slide.placeholders[1].text = (
    "Imprecise probabilities help manage uncertainty in segmentation.\n"
    "Calibration and uncertainty metrics are crucial for reliable models.\n"
    "Future work: Explore more uncertainty metrics and applications.\n"
)

# Save the presentation
prs.save("imprecise_probability_presentation/slides/imprecise_probabilityV2_summary.pptx")